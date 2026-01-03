# sensitive_data_generator/advanced_file_writers.py

from __future__ import annotations

import os
import random
from datetime import datetime, timedelta
import pandas as pd
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Paragraph, Spacer
from reportlab.lib.units import inch
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import numpy as np
import io
from .generators import PIIGenerator

class AdvancedFileWriter:
    """Advanced multi-format file writers for synthetic datasets."""

    @staticmethod
    def create_complex_pdf(content, output_dir, filename=None, include_charts=True):
        """Create a complex PDF document (tables + charts)."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"complex_document_{timestamp}.pdf"

        filepath = os.path.join(output_dir, filename)

        # Build the PDF document.
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        elements = []
        styles = getSampleStyleSheet()

        # Add title.
        title_style = ParagraphStyle(
            'Title',
            parent=styles['Heading1'],
            fontSize=18,
            alignment=1,  # centered
            spaceAfter=12
        )
        title = Paragraph("機密文件 - 個人資料報告", title_style)
        elements.append(title)

        # Add content paragraph.
        pii_style = ParagraphStyle(
            'PIIContent',
            parent=styles['BodyText'],
            fontSize=10,
            leading=14,
            spaceAfter=6
        )
        pii_paragraph = Paragraph(content, pii_style)
        elements.append(pii_paragraph)
        elements.append(Spacer(1, 12))

        # Add a table.
        table_data = [
            ['項目', '原始資料', '備註'],
            ['姓名', PIIGenerator.generate_tw_name(), '測試用虛構姓名'],
            ['身分證字號', PIIGenerator.generate_tw_id(), '測試用虛構ID'],
            ['電話', PIIGenerator.generate_tw_phone(), '測試用虛構電話'],
            ['地址', PIIGenerator.generate_tw_address(), '測試用虛構地址']
        ]

        table = Table(table_data, colWidths=[1.5*inch, 3*inch, 2.5*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        elements.append(table)
        elements.append(Spacer(1, 24))

        # Add a chart.
        if include_charts:
            chart_img = AdvancedFileWriter.generate_fake_chart()
            elements.append(RLImage(chart_img, width=5*inch, height=3*inch))
            elements.append(Spacer(1, 12))
            chart_caption = Paragraph("圖1: 測試資料分布圖", styles['Italic'])
            elements.append(chart_caption)

        # Write PDF.
        doc.build(elements)
        return filepath

    @staticmethod
    def generate_fake_chart():
        """Generate a fake chart image for reports."""
        plt.figure(figsize=(8, 5))

        # Randomly choose a chart type.
        chart_type = random.choice(['bar', 'line', 'pie'])

        if chart_type == 'bar':
            labels = ['A部門', 'B部門', 'C部門', 'D部門']
            values = [random.randint(100, 500) for _ in range(4)]
            plt.bar(labels, values, color=plt.cm.Pastel1(range(4)))
            plt.title('部門業績比較')
            plt.ylabel('業績 (萬元)')
        elif chart_type == 'line':
            x = np.arange(1, 11)
            y = np.random.rand(10) * 100
            plt.plot(x, y, marker='o', linestyle='-')
            plt.title('月度趨勢分析')
            plt.xlabel('月份')
            plt.ylabel('數值')
            plt.xticks(x, [f'{i}月' for i in range(1, 11)])
        else:  # pie
            labels = ['類別A', '類別B', '類別C', '類別D']
            sizes = [random.randint(15, 40) for _ in range(4)]
            plt.pie(sizes, labels=labels, autopct='%1.1f%%', shadow=True)
            plt.title('類別分布圖')

        plt.tight_layout()

        # Save to BytesIO.
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=100)
        plt.close()
        img_buffer.seek(0)
        return img_buffer

    @staticmethod
    def create_word_document(content, output_dir, filename=None):
        """Create a complex Word document (tables + images)."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"complex_document_{timestamp}.docx"

        filepath = os.path.join(output_dir, filename)

        # Build the Word document.
        doc = Document()

        # Add title.
        title = doc.add_heading('機密文件 - 個人資料報告', level=0)
        title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Add date.
        date_para = doc.add_paragraph()
        date_para.add_run(f"生成日期: {datetime.now().strftime('%Y-%m-%d')}").italic = True
        date_para.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT

        # Add content paragraph.
        doc.add_paragraph("以下為測試用虛構個人資料:")
        pii_para = doc.add_paragraph(content)

        # Add a table.
        table = doc.add_table(rows=5, cols=3)
        table.style = 'LightShading-Accent1'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER

        # Header row.
        header_cells = table.rows[0].cells
        header_cells[0].text = '項目'
        header_cells[1].text = '原始資料'
        header_cells[2].text = '備註'

        # Table rows.
        data_rows = [
            ('姓名', PIIGenerator.generate_tw_name(), '測試用虛構姓名'),
            ('身分證字號', PIIGenerator.generate_tw_id(), '測試用虛構ID'),
            ('電話', PIIGenerator.generate_tw_phone(), '測試用虛構電話'),
            ('地址', PIIGenerator.generate_tw_address(), '測試用虛構地址')
        ]

        for i, row_data in enumerate(data_rows, start=1):
            row_cells = table.rows[i].cells
            row_cells[0].text = row_data[0]
            row_cells[1].text = row_data[1]
            row_cells[2].text = row_data[2]

        # Add a chart image.
        doc.add_paragraph("\n資料分布圖表:")
        chart_img = AdvancedFileWriter.generate_fake_chart()
        chart_img.seek(0)
        doc.add_picture(chart_img, width=Inches(5.0))

        # Add footer.
        section = doc.sections[0]
        footer = section.footer
        footer_para = footer.paragraphs[0]
        footer_para.text = "本文件包含測試用虛構個人資料 - 機密文件"
        footer_para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

        # Save document.
        doc.save(filepath)
        return filepath

    @staticmethod
    def create_excel_spreadsheet(output_dir, filename=None):
        """Create an Excel spreadsheet containing synthetic sensitive data."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"financial_data_{timestamp}.xlsx"

        filepath = os.path.join(output_dir, filename)

        # Build a dataframe.
        data = {
            '客戶編號': [f"CUST-{random.randint(1000, 9999)}" for _ in range(20)],
            '客戶姓名': [PIIGenerator.generate_tw_name() for _ in range(20)],
            '身分證字號': [PIIGenerator.generate_tw_id() for _ in range(20)],
            '電話號碼': [PIIGenerator.generate_tw_phone() for _ in range(20)],
            '電子郵件': [PIIGenerator.generate_email() for _ in range(20)],
            '帳戶餘額': [round(random.uniform(1000, 100000), 2) for _ in range(20)],
            '最近交易日期': [
                (datetime.now() - timedelta(days=random.randint(1, 365))).strftime("%Y-%m-%d")
                for _ in range(20)
            ],
        }

        df = pd.DataFrame(data)

        # Write via xlsxwriter.
        with pd.ExcelWriter(filepath, engine='xlsxwriter') as writer:
            df.to_excel(writer, sheet_name='客戶資料', index=False)

            # Get workbook and worksheet objects.
            workbook = writer.book
            worksheet = writer.sheets['客戶資料']

            # Define a header format.
            header_format = workbook.add_format({
                'bold': True,
                'text_wrap': True,
                'valign': 'top',
                'fg_color': '#D7E4BC',
                'border': 1
            })

            # Apply header format.
            for col_num, value in enumerate(df.columns.values):
                worksheet.write(0, col_num, value, header_format)

            # Set column widths.
            worksheet.set_column('A:A', 12)
            worksheet.set_column('B:B', 15)
            worksheet.set_column('C:C', 15)
            worksheet.set_column('D:D', 15)
            worksheet.set_column('E:E', 25)
            worksheet.set_column('F:F', 15)
            worksheet.set_column('G:G', 15)

            # Add a chart.
            chart = workbook.add_chart({'type': 'column'})
            chart.add_series({
                'name': '=客戶資料!$F$1',
                'categories': '=客戶資料!$B$2:$B$21',
                'values': '=客戶資料!$F$2:$F$21',
            })
            chart.set_title({'name': '客戶帳戶餘額分布'})
            chart.set_x_axis({'name': '客戶姓名'})
            chart.set_y_axis({'name': '餘額 (NT$)'})
            worksheet.insert_chart('I2', chart)

        return filepath

    @staticmethod
    def create_powerpoint_presentation(output_dir, filename=None):
        """Create a PowerPoint presentation containing synthetic sensitive data."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"business_report_{timestamp}.pptx"

        filepath = os.path.join(output_dir, filename)

        # Build presentation.
        prs = Presentation()

        # Title slide.
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "業務分析報告"
        subtitle.text = f"報告日期: {datetime.now().strftime('%Y-%m-%d')}\n機密文件"

        # Content slide: text + table.
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "客戶資料摘要"

        content = slide.placeholders[1]
        content.text = (
            "主要客戶統計:\n"
            f"- 客戶總數: 20\n"
            f"- 平均帳戶餘額: NT$ {random.randint(50000, 100000):,}\n"
            f"- 最近新增客戶: {PIIGenerator.generate_tw_name()}\n\n"
            "以下為抽樣資料:"
        )

        # Add a table.
        rows, cols = 4, 3
        left = Inches(1.0)
        top = Inches(2.5)
        width = Inches(8.0)
        height = Inches(1.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row.
        table.cell(0, 0).text = '客戶姓名'
        table.cell(0, 1).text = '電話號碼'
        table.cell(0, 2).text = '帳戶餘額'

        # Fill table rows.
        for i in range(1, rows):
            table.cell(i, 0).text = PIIGenerator.generate_tw_name()
            table.cell(i, 1).text = PIIGenerator.generate_tw_phone()
            table.cell(i, 2).text = f"NT$ {random.randint(10000, 100000):,}"

        # Chart slide.
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "業績分析圖表"

        # Add a chart image.
        chart_img = AdvancedFileWriter.generate_fake_chart()
        chart_img.seek(0)
        left = Inches(1)
        top = Inches(1.5)
        slide.shapes.add_picture(chart_img, left, top, width=Inches(8))

        # Save presentation.
        prs.save(filepath)
        return filepath

    @staticmethod
    def create_scanned_document(content, output_dir, filename=None):
        """Create a scanned-document style image (stamp + signature)."""
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"scanned_doc_{timestamp}.png"

        filepath = os.path.join(output_dir, filename)

        # Create background (paper simulation).
        img = Image.new('RGB', (1240, 1754), color=(248, 246, 240))
        draw = ImageDraw.Draw(img)

        # Add paper texture.
        for _ in range(500):
            x = random.randint(0, 1240)
            y = random.randint(0, 1754)
            r = random.randint(1, 3)
            draw.ellipse([x, y, x+r, y+r], fill=(220, 220, 220))

        # Use Traditional Chinese fonts if available.
        try:
            font = ImageFont.truetype("fonts/NotoSansTC-Regular.otf", 24)
            title_font = ImageFont.truetype("fonts/NotoSansTC-Bold.otf", 36)
        except:
            font = ImageFont.load_default()
            title_font = ImageFont.load_default()

        # Add title.
        draw.text((620, 100), "個人資料表", fill=(0, 0, 0), font=title_font, anchor="mm")

        # Add content.
        y_position = 200
        for line in content.split('\n'):
            draw.text((100, y_position), line, fill=(0, 0, 0), font=font)
            y_position += 40
            if y_position > 1600:
                break

        # Add a table.
        draw.rectangle([80, 300, 1160, 800], outline=(0, 0, 0), width=2)

        # Horizontal lines.
        for y in [380, 460, 540, 620, 700]:
            draw.line([80, y, 1160, y], fill=(0, 0, 0), width=1)

        # Vertical lines.
        draw.line([300, 300, 300, 800], fill=(0, 0, 0), width=1)
        draw.line([700, 300, 700, 800], fill=(0, 0, 0), width=1)

        # Table header.
        headers = ["項目", "資料", "備註"]
        draw.text((190, 330), headers[0], fill=(0, 0, 0), font=font, anchor="mm")
        draw.text((500, 330), headers[1], fill=(0, 0, 0), font=font, anchor="mm")
        draw.text((930, 330), headers[2], fill=(0, 0, 0), font=font, anchor="mm")

        # Table rows.
        rows = [
            ("姓名", PIIGenerator.generate_tw_name(), "測試用虛構姓名"),
            ("身分證字號", PIIGenerator.generate_tw_id(), "測試用虛構ID"),
            ("電話號碼", PIIGenerator.generate_tw_phone(), "測試用虛構電話"),
            ("地址", PIIGenerator.generate_tw_address(), "測試用虛構地址")
        ]

        for i, row in enumerate(rows):
            y_pos = 330 + 80 * (i+1)
            draw.text((190, y_pos), row[0], fill=(0, 0, 0), font=font, anchor="mm")
            draw.text((500, y_pos), row[1], fill=(0, 0, 0), font=font, anchor="mm")
            draw.text((930, y_pos), row[2], fill=(0, 0, 0), font=font, anchor="mm")

        # Add a stamp.
        stamp_size = 150
        stamp_x = 1000
        stamp_y = 900
        draw.ellipse([stamp_x, stamp_y, stamp_x+stamp_size, stamp_y+stamp_size],
                     outline=(200, 0, 0), width=3)

        # Stamp text.
        stamp_font = ImageFont.truetype("fonts/NotoSansTC-Bold.otf", 20)
        draw.text((stamp_x+stamp_size/2, stamp_y+stamp_size/2-15), "核准章",
                 fill=(200, 0, 0), font=stamp_font, anchor="mm")
        draw.text((stamp_x+stamp_size/2, stamp_y+stamp_size/2+15), "測試專用",
                 fill=(200, 0, 0), font=stamp_font, anchor="mm")

        # Add a signature.
        signature_x = 200
        signature_y = 900
        draw.line([signature_x, signature_y, signature_x+150, signature_y-50],
                 fill=(0, 0, 0), width=3)
        draw.line([signature_x+150, signature_y-50, signature_x+200, signature_y+50],
                 fill=(0, 0, 0), width=3)
        draw.text((signature_x+100, signature_y+70), PIIGenerator.generate_tw_name(),
                 fill=(0, 0, 0), font=font)

        # Save image.
        img.save(filepath)
        return filepath
