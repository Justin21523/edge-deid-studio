from __future__ import annotations

from pathlib import Path

import pytest

from deid_pipeline.config import Config
from deid_pipeline.core.anchors import attach_segment_anchors
from deid_pipeline.core.contracts import normalize_entity, replacement_key
from deid_pipeline.pii.detectors.regex_detector import RegexDetector


def _detect_entities(text: str, *, language: str) -> list[dict]:
    cfg = Config()
    rules = cfg.REGEX_RULES_FILE if language == "zh" else cfg.REGEX_EN_RULES_FILE
    detector = RegexDetector(rules)
    return detector.detect(text)


def test_docx_rebuild_writes_output_file(tmp_path: Path) -> None:
    pytest.importorskip("docx")
    from docx import Document  # type: ignore

    source = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Contact: alice@example.com")
    doc.save(str(source))

    from deid_pipeline.handlers.docx import DocxHandler

    handler = DocxHandler()
    extracted = handler.extract(source, language="en")
    raw_entities = _detect_entities(extracted.text, language="en")
    entities = [normalize_entity(e, language="en", text=extracted.text) for e in raw_entities]
    attach_segment_anchors(entities, extracted.segments)

    replacement_map = {replacement_key("EMAIL", "alice@example.com"): "<EMAIL_REDACTED>"}
    artifacts = handler.rebuild(
        extracted,
        output_text=extracted.text.replace("alice@example.com", "<EMAIL_REDACTED>"),
        entities=entities,
        replacement_map=replacement_map,
        events=[],
        output_dir=tmp_path,
        mode="replace",
    )

    assert artifacts.get("rebuild_supported") is True
    out_path = Path(str(artifacts["output_path"]))
    assert out_path.exists()

    out_doc = Document(str(out_path))
    out_text = "\n".join(p.text for p in out_doc.paragraphs)
    assert "alice@example.com" not in out_text
    assert "<EMAIL_REDACTED>" in out_text


def test_xlsx_rebuild_writes_output_file(tmp_path: Path) -> None:
    pytest.importorskip("openpyxl")
    from openpyxl import Workbook, load_workbook  # type: ignore

    source = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"].value = "alice@example.com"
    wb.save(str(source))
    wb.close()

    from deid_pipeline.handlers.xlsx import XlsxHandler

    handler = XlsxHandler()
    extracted = handler.extract(source, language="en")
    raw_entities = _detect_entities(extracted.text, language="en")
    entities = [normalize_entity(e, language="en", text=extracted.text) for e in raw_entities]
    attach_segment_anchors(entities, extracted.segments)

    replacement_map = {replacement_key("EMAIL", "alice@example.com"): "<EMAIL_REDACTED>"}
    artifacts = handler.rebuild(
        extracted,
        output_text=extracted.text.replace("alice@example.com", "<EMAIL_REDACTED>"),
        entities=entities,
        replacement_map=replacement_map,
        events=[],
        output_dir=tmp_path,
        mode="replace",
    )

    assert artifacts.get("rebuild_supported") is True
    out_path = Path(str(artifacts["output_path"]))
    assert out_path.exists()

    out_wb = load_workbook(filename=str(out_path), data_only=True)
    out_ws = out_wb.active
    assert out_ws["A1"].value == "<EMAIL_REDACTED>"
    out_wb.close()


def test_pptx_rebuild_writes_output_file(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    source = tmp_path / "sample.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])  # blank
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.text_frame.text = "Contact: alice@example.com"
    pres.save(str(source))

    from deid_pipeline.handlers.pptx import PptxHandler

    handler = PptxHandler()
    extracted = handler.extract(source, language="en")
    raw_entities = _detect_entities(extracted.text, language="en")
    entities = [normalize_entity(e, language="en", text=extracted.text) for e in raw_entities]
    attach_segment_anchors(entities, extracted.segments)

    replacement_map = {replacement_key("EMAIL", "alice@example.com"): "<EMAIL_REDACTED>"}
    artifacts = handler.rebuild(
        extracted,
        output_text=extracted.text.replace("alice@example.com", "<EMAIL_REDACTED>"),
        entities=entities,
        replacement_map=replacement_map,
        events=[],
        output_dir=tmp_path,
        mode="replace",
    )

    assert artifacts.get("rebuild_supported") is True
    out_path = Path(str(artifacts["output_path"]))
    assert out_path.exists()

    out_pres = Presentation(str(out_path))
    texts: list[str] = []
    for slide in out_pres.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)

    combined = "\n".join(texts)
    assert "alice@example.com" not in combined
    assert "<EMAIL_REDACTED>" in combined

