from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Optional

from .base import ExtractedDocument, SegmentSpec, build_document_from_segments
from ..core.contracts import replacement_key


class PptxHandler:
    extensions = [".pptx"]

    def extract(self, input_path: Path, *, language: str) -> ExtractedDocument:
        try:
            from pptx import Presentation  # type: ignore
        except Exception as exc:
            raise ImportError("python-pptx is required to extract PPTX files") from exc

        pres = Presentation(str(input_path))
        segments: list[SegmentSpec] = []

        for slide_idx, slide in enumerate(pres.slides):
            for shape in slide.shapes:
                text = getattr(shape, "text", "") or ""
                text = str(text).strip()
                if not text:
                    continue
                segments.append(
                    SegmentSpec(
                        text=text,
                        metadata={
                            "kind": "shape_text",
                            "slide_index": slide_idx,
                            "shape_id": getattr(shape, "shape_id", None),
                        },
                    )
                )

        return build_document_from_segments(
            input_path=input_path,
            language=language,
            segments=segments or [SegmentSpec(text="")],
            separator="\n",
            metadata={"format": "pptx", "slides": len(pres.slides)},
        )

    def rebuild(
        self,
        document: ExtractedDocument,
        *,
        output_text: str,
        entities: list[dict],
        replacement_map: Dict[str, str],
        events: list[dict],
        output_dir: Optional[Path] = None,
        mode: str = "replace",
    ) -> Dict[str, Any]:
        artifacts: Dict[str, Any] = {"output_text": output_text}

        if output_dir is None or mode != "replace":
            artifacts["rebuild_supported"] = False
            return artifacts

        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            artifacts["rebuild_supported"] = False
            return artifacts

        output_dir.mkdir(parents=True, exist_ok=True)
        out_path = output_dir / f"{document.input_path.stem}.deid.pptx"

        pres = Presentation(str(document.input_path))

        replacements: list[tuple[str, str]] = []
        for entity in entities:
            original = entity.get("text")
            entity_type = entity.get("type")
            if not original or not entity_type:
                continue
            repl = replacement_map.get(replacement_key(str(entity_type), str(original)))
            if repl is None:
                continue
            replacements.append((str(original), str(repl)))

        for slide in pres.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = getattr(shape, "text", "") or ""
                new_text = str(text)
                for original, repl in replacements:
                    new_text = new_text.replace(original, repl)
                if new_text != text:
                    shape.text = new_text

        pres.save(str(out_path))

        artifacts["output_path"] = str(out_path)
        artifacts["rebuild_supported"] = True
        artifacts["rebuild_mode"] = "pptx_text_replace"
        return artifacts
