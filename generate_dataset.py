import os
import json
import time
from datetime import datetime
import matplotlib
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from PIL import ImageFont
import openpyxl
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from sensitive_data_generator.dataset_generator import MultiFormatDatasetGenerator
from sensitive_data_generator.advanced_file_writers import AdvancedFileWriter
from sensitive_data_generator.formatters import DataFormatter
from sensitive_data_generator.advanced_formatters import AdvancedDataFormatter
from sensitive_data_generator.generators import 

# 取得專案根目錄
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

# 1) 註冊 ReportLab PDF 用中文字體
font_path = os.path.join(PROJECT_ROOT , "sensitive_data_generator", "fonts", "NotoSansTC-Regular.ttf")
pdfmetrics.registerFont(TTFont("NotoSansTC", font_path))

# 2) 設定 matplotlib 的中文字型（畫圖）
matplotlib.font_manager.fontManager.addfont(font_path)
matplotlib.rcParams["font.family"] = "NotoSansTC"

# 3) PIL 如果有用到 ImageFont，也可以直接共用
PIL_FONT_REGULAR = ImageFont.truetype(font_path, size=24)
PIL_FONT_BOLD    = ImageFont.truetype(
    os.path.join(PROJECT_ROOT , "sensitive_data_generator", "fonts", "NotoSansTC-Bold.ttf"),
    size=36
)

# 以下是範例

generator = DatasetGenerator(
    document_type="medical_record",
    format="pdf",
    count=50,
    locale="zh_TW"
)
generator.generate("output/medical_records")

def main():
    # 設定輸出目錄和數量
    OUTPUT_DIR = "generated_dataset"
    NUM_ITEMS = 50  # 生成50個完整項目（每個項目包含多種格式）

    print(f"開始生成測試資料集，數量: {NUM_ITEMS}...")
    start_time = time.time()

    # 步驟1: 生成完整多格式資料集
    dataset = MultiFormatDatasetGenerator.generate_full_dataset(
        output_dir=OUTPUT_DIR,
        num_items=NUM_ITEMS
    )

    # 步驟2: 額外生成專業文件範例
    generate_professional_examples(OUTPUT_DIR)

    # 步驟3: 生成資料集摘要報告
    generate_summary_report(dataset, OUTPUT_DIR)

    duration = time.time() - start_time
    print(f"資料集生成完成! 共生成 {NUM_ITEMS} 個項目，耗時: {duration:.2f} 秒")
    print(f"輸出目錄: {os.path.abspath(OUTPUT_DIR)}")

def generate_professional_examples(output_dir):
    """生成多樣化的專業文件範例（PDF, DOCX, XLSX, CSV, PPTX, PNG, TXT)"""
    examples_dir = os.path.join(output_dir, "examples")
    os.makedirs(examples_dir, exist_ok=True)

    # --- 1) 自然語言段落 範例 (TXT + CSV) ---
    paragraph = DataFormatter.generate_paragraph()  # 自訂方法：回傳一段含 PII 的文字
    txt_path = os.path.join(examples_dir, "sample_paragraph.txt")
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(paragraph)

    df = pd.DataFrame([{"text": line} for line in paragraph.split("\n")])
    csv_path = os.path.join(examples_dir, "sample_paragraph.csv")
    df.to_csv(csv_path, index=False, encoding="utf-8-sig")

    # --- 2) 合約文件 範例 (PDF + PPTX) ---
    contract_content = AdvancedDataFormatter.generate_contract_document()
    # PDF
    AdvancedFileWriter.create_complex_pdf(
        contract_content,
        examples_dir,
        "professional_contract.pdf",
        font_name="NotoSansTC",       # ReportLab 已註冊的 font
        font_size=12
    )
    # PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(800), Pt(400)).text_frame
    p = tx.add_paragraph()
    p.text = paragraph  # 或 contract_content 裡的某一段
    p.font.name = "Noto Sans TC"
    p.font.size = Pt(18)
    prs.save(os.path.join(examples_dir, "professional_contract.pptx"))

    # --- 3) 醫療報告 範例 (DOCX + XLSX + PNG) ---
    medical_content = AdvancedDataFormatter.generate_medical_report()
    # DOCX
    AdvancedFileWriter.create_word_document(
        medical_content,
        examples_dir,
        "professional_medical_report.docx",
        font_name="Noto Sans TC",
        font_size=12
    )
    # XLSX
    AdvancedFileWriter.create_excel_spreadsheet(
        examples_dir,
        "professional_medical_data.xlsx",
        data=medical_content["table_data"],  # 假設是一個 list of dict
        font_name="Noto Sans TC"
    )
    # PNG（示意：把醫療資料畫成長條圖）
    ages = [row["age"] for row in medical_content["table_data"]]
    counts = [row["count"] for row in medical_content["table_data"]]
    plt.figure(figsize=(6,4))
    plt.bar(ages, counts)
    plt.title("年齡分佈")
    plt.xlabel("年齡")
    plt.ylabel("人數")
    png_path = os.path.join(examples_dir, "medical_chart.png")
    plt.savefig(png_path, bbox_inches="tight")

    # --- 4) 財務報表 範例 (XLSX + PNG) ---
    fin_data = AdvancedDataFormatter.generate_financial_statement()
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["日期", "描述", "金額"])
    for tx in fin_data:
        ws.append([tx["date"], tx["description"], tx["amount"]])
    for row in ws.rows:
        for cell in row:
            cell.font = openpyxl.styles.Font(name="Noto Sans TC")
    xlsx_path = os.path.join(examples_dir, "professional_financial_report.xlsx")
    wb.save(xlsx_path)

    # 財務圖表
    dates = [tx["date"] for tx in fin_data]
    amounts = [tx["amount"] for tx in fin_data]
    plt.figure(figsize=(6,4))
    plt.plot(dates, amounts, marker="o")
    plt.title("交易金額走勢")
    plt.xlabel("日期")
    plt.ylabel("金額 (元)")
    plt.xticks(rotation=45)
    png2 = os.path.join(examples_dir, "financial_chart.png")
    plt.savefig(png2, bbox_inches="tight")

    # 生成掃描文件範例
    AdvancedFileWriter.create_scanned_document(
        contract_content,
        examples_dir,
        "scanned_contract_example.png"
    )

def generate_summary_report(dataset, output_dir):
    """生成資料集摘要報告"""
    report = {
        "generated_at": datetime.now().isoformat(),
        "total_items": len(dataset),
        "format_distribution": {
            "pdf": 0,
            "word": 0,
            "excel": 0,
            "ppt": 0,
            "scanned": 0,
            "contract": 0,
            "medical": 0,
            "financial": 0
        },
        "pii_type_distribution": {
            "TW_ID": 0,
            "PHONE": 0,
            "ADDRESS": 0,
            "NAME": 0,
            "MEDICAL_RECORD": 0,
            "DATE_OF_BIRTH": 0,
            "EMAIL": 0,
            "CREDIT_CARD": 0,
            "PASSPORT": 0,
            "LICENSE_PLATE": 0,
            "HEALTH_INSURANCE": 0
        },
        "sample_items": []
    }

    # 分析資料集
    for item in dataset:
        # 統計文件類型
        doc_type = item["metadata"]["type"]
        report["format_distribution"][doc_type] += 1

        # 統計格式
        for file_info in item["formats"]:
            fmt = file_info["format"]
            if fmt in report["format_distribution"]:
                report["format_distribution"][fmt] += 1

        # 提取前3個項目作為樣本
        if len(report["sample_items"]) < 3:
            sample_item = {
                "id": item["id"],
                "type": doc_type,
                "formats": [f["format"] for f in item["formats"]],
                "content_snippet": item["metadata"]["content"][:200] + "..."
            }
            report["sample_items"].append(sample_item)

    # 保存報告
    report_path = os.path.join(output_dir, "dataset_summary.json")
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)

    # 輸出統計摘要
    print("\n資料集統計摘要:")
    print(f"總項目數: {report['total_items']}")
    print("文件類型分佈:")
    for doc_type, count in report["format_distribution"].items():
        if count > 0:
            print(f"  {doc_type}: {count}")

if __name__ == "__main__":
    main()
